#!/usr/bin/env python3
"""PPTX to PDF converter using python-pptx + weasyprint.

Replaces LibreOffice's `soffice --headless --convert-to pdf` for presentations.

Pipeline: PPTX → python-pptx (parse) → HTML/CSS per slide → weasyprint → PDF

Usage:
    python pptx_to_pdf.py <input.pptx> <output.pdf>

Dependencies (pip):
    python-pptx, weasyprint, Pillow
"""

import sys
import os
import io
import base64
import tempfile
from pathlib import Path


def emu_to_pt(emu):
    """Convert EMU (English Metric Units) to CSS points."""
    return emu / 12700.0


def emu_to_in(emu):
    """Convert EMU to inches."""
    return emu / 914400.0


def color_to_css(color):
    """Convert a python-pptx color object to a CSS color string."""
    try:
        if color and color.type is not None:
            rgb = color.rgb
            if rgb:
                return f"#{rgb}"
    except (AttributeError, TypeError):
        pass
    return None


def fill_to_css(fill):
    """Convert a python-pptx FillFormat to CSS background properties."""
    try:
        fill_type = fill.type
    except (AttributeError, TypeError):
        return ""

    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.dml.color import RGBColor

    if fill_type is None:
        return ""

    # Solid fill
    try:
        if fill.type is not None:
            fc = fill.fore_color
            if fc and fc.type is not None:
                rgb = fc.rgb
                if rgb:
                    return f"background-color: #{rgb};"
    except (AttributeError, TypeError, KeyError):
        pass

    return ""


def render_text_frame(tf):
    """Render a TextFrame to HTML."""
    html_parts = []
    for para in tf.paragraphs:
        parts = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            # Escape HTML entities
            text = (
                text.replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")
            )

            style = []
            font = run.font
            if font.bold:
                style.append("font-weight:bold")
            if font.italic:
                style.append("font-style:italic")
            if font.underline:
                style.append("text-decoration:underline")
            if font.size:
                style.append(f"font-size:{emu_to_pt(font.size):.1f}pt")
            if font.name:
                style.append(f"font-family:'{font.name}',sans-serif")

            fc = color_to_css(font.color)
            if fc:
                style.append(f"color:{fc}")

            if style:
                parts.append(f'<span style="{";".join(style)}">{text}</span>')
            else:
                parts.append(text)

        # Paragraph alignment
        from pptx.enum.text import PP_ALIGN

        p_style = []
        if para.alignment == PP_ALIGN.CENTER:
            p_style.append("text-align:center")
        elif para.alignment == PP_ALIGN.RIGHT:
            p_style.append("text-align:right")
        elif para.alignment == PP_ALIGN.JUSTIFY:
            p_style.append("text-align:justify")

        if para.space_before:
            p_style.append(f"margin-top:{emu_to_pt(para.space_before):.1f}pt")
        if para.space_after:
            p_style.append(f"margin-bottom:{emu_to_pt(para.space_after):.1f}pt")

        style_attr = f' style="{";".join(p_style)}"' if p_style else ""
        content = "".join(parts) if parts else "&nbsp;"
        html_parts.append(f"<p{style_attr}>{content}</p>")

    return "\n".join(html_parts)


def render_table(table):
    """Render a Table shape to HTML."""
    rows_html = []
    for row in table.rows:
        cells_html = []
        for cell in row.cells:
            # Get cell text with basic formatting
            cell_text = render_text_frame(cell.text_frame)
            style = []

            # Cell fill
            try:
                fc = cell.fill
                if fc.type is not None:
                    rgb = fc.fore_color.rgb
                    if rgb:
                        style.append(f"background-color:#{rgb}")
            except (AttributeError, TypeError, KeyError):
                pass

            style.append("border:1px solid #999")
            style.append("padding:4pt 6pt")
            style.append("vertical-align:top")

            style_attr = f' style="{";".join(style)}"'
            cells_html.append(f"<td{style_attr}>{cell_text}</td>")
        rows_html.append(f"<tr>{''.join(cells_html)}</tr>")

    return f'<table style="border-collapse:collapse;position:absolute">{"".join(rows_html)}</table>'


def extract_image_data(image):
    """Extract image data as a base64 data URI."""
    try:
        blob = image.blob
        content_type = image.content_type
        b64 = base64.b64encode(blob).decode("ascii")
        return f"data:{content_type};base64,{b64}"
    except Exception:
        return None


def pptx_to_pdf(input_path: str, output_path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from weasyprint import HTML

    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)

    if not os.path.exists(input_path):
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(input_path)

    # Get slide dimensions
    slide_w_in = emu_to_in(prs.slide_width)
    slide_h_in = emu_to_in(prs.slide_height)

    slides_html = []

    for slide_idx, slide in enumerate(prs.slides):
        shapes_html = []

        # Slide background
        bg_css = ""
        try:
            bg = slide.background
            fill = bg.fill
            bg_css = fill_to_css(fill)
        except (AttributeError, TypeError):
            pass

        if not bg_css:
            # Try slide layout background
            try:
                bg = slide.slide_layout.background
                fill = bg.fill
                bg_css = fill_to_css(fill)
            except (AttributeError, TypeError):
                pass

        if not bg_css:
            bg_css = "background-color: white;"

        for shape in slide.shapes:
            left = emu_to_pt(shape.left) if shape.left else 0
            top = emu_to_pt(shape.top) if shape.top else 0
            width = emu_to_pt(shape.width) if shape.width else 0
            height = emu_to_pt(shape.height) if shape.height else 0

            shape_style = (
                f"position:absolute;"
                f"left:{left:.1f}pt;"
                f"top:{top:.1f}pt;"
                f"width:{width:.1f}pt;"
                f"height:{height:.1f}pt;"
                f"overflow:hidden;"
            )

            # Rotation
            if shape.rotation:
                shape_style += f"transform:rotate({shape.rotation}deg);"

            # Shape fill
            try:
                if hasattr(shape, "fill"):
                    shape_fill = fill_to_css(shape.fill)
                    if shape_fill:
                        shape_style += shape_fill
            except (AttributeError, TypeError):
                pass

            # Image shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                data_uri = extract_image_data(shape.image)
                if data_uri:
                    shapes_html.append(
                        f'<div style="{shape_style}">'
                        f'<img src="{data_uri}" '
                        f'style="width:100%;height:100%;object-fit:contain;">'
                        f"</div>"
                    )
                continue

            # Table shapes
            if shape.has_table:
                table_html = render_table(shape.table)
                shapes_html.append(
                    f'<div style="{shape_style}">{table_html}</div>'
                )
                continue

            # Group shapes — render child shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Flatten group: each child shape is positioned relative to group
                try:
                    for child in shape.shapes:
                        cl = emu_to_pt(child.left) if child.left else 0
                        ct = emu_to_pt(child.top) if child.top else 0
                        cw = emu_to_pt(child.width) if child.width else 0
                        ch = emu_to_pt(child.height) if child.height else 0
                        child_style = (
                            f"position:absolute;"
                            f"left:{cl:.1f}pt;top:{ct:.1f}pt;"
                            f"width:{cw:.1f}pt;height:{ch:.1f}pt;"
                            f"overflow:hidden;"
                        )
                        if hasattr(child, "text_frame"):
                            content = render_text_frame(child.text_frame)
                            shapes_html.append(
                                f'<div style="{child_style}">{content}</div>'
                            )
                except (AttributeError, TypeError):
                    pass
                continue

            # Text shapes (most common)
            if shape.has_text_frame:
                content = render_text_frame(shape.text_frame)
                shapes_html.append(
                    f'<div style="{shape_style}">{content}</div>'
                )
                continue

            # Placeholder shapes with text
            try:
                if hasattr(shape, "text") and shape.text:
                    text = (
                        shape.text.replace("&", "&amp;")
                        .replace("<", "&lt;")
                        .replace(">", "&gt;")
                    )
                    shapes_html.append(
                        f'<div style="{shape_style}"><p>{text}</p></div>'
                    )
            except (AttributeError, TypeError):
                pass

        slide_html = (
            f'<div class="slide" style="{bg_css}">'
            f'{"".join(shapes_html)}'
            f"</div>"
        )
        slides_html.append(slide_html)

    # Build full HTML document
    full_html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
@page {{
    size: {slide_w_in:.3f}in {slide_h_in:.3f}in;
    margin: 0;
}}
body {{
    margin: 0;
    padding: 0;
}}
.slide {{
    position: relative;
    width: {emu_to_pt(prs.slide_width):.1f}pt;
    height: {emu_to_pt(prs.slide_height):.1f}pt;
    overflow: hidden;
    page-break-after: always;
    box-sizing: border-box;
}}
.slide:last-child {{
    page-break-after: avoid;
}}
p {{
    margin: 0;
    padding: 0;
    line-height: 1.2;
}}
table {{
    border-collapse: collapse;
}}
</style>
</head>
<body>
{''.join(slides_html)}
</body>
</html>"""

    # Render to PDF
    HTML(string=full_html).write_pdf(output_path)

    if not os.path.exists(output_path):
        print("Error: PDF generation failed", file=sys.stderr)
        sys.exit(1)

    size = os.path.getsize(output_path)
    print(f"OK:{size}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} <input.pptx> <output.pdf>", file=sys.stderr)
        sys.exit(1)
    pptx_to_pdf(sys.argv[1], sys.argv[2])
